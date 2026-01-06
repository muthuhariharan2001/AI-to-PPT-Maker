import os
from rest_framework.decorators import api_view
from rest_framework.response import Response
from django.conf import settings
from pptx import Presentation
import google.generativeai as genai

# ✅ Configure Gemini API
genai.configure(api_key="AIzaSyARicFLSRPml019dtya6Y27EVPCBfHMDdI")

# ✅ Initialize with proper model name
model = genai.GenerativeModel(model_name="models/gemini-1.5-pro")

def generate_gemini_slide_content(topic, slide_number):
    prompt = f"Create 3 simple bullet points for a PowerPoint slide about '{topic}' - Slide {slide_number}."
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        return f"Error generating content: {str(e)}"

@api_view(['POST'])
def generate_ppt(request):
    topic = request.data.get("topic")
    slide_count = int(request.data.get("slideCount", 5))

    if not topic:
        return Response({"error": "Topic is required."}, status=400)

    prs = Presentation()

    for i in range(slide_count):
        content = generate_gemini_slide_content(topic, i + 1)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{topic} - Slide {i+1}"
        slide.placeholders[1].text = content

    # Save PPT
    media_dir = settings.MEDIA_ROOT
    os.makedirs(media_dir, exist_ok=True)

    filename = f"{topic.replace(' ', '_')}.pptx"
    save_path = os.path.join(media_dir, filename)
    prs.save(save_path)

    file_url = request.build_absolute_uri(f"{settings.MEDIA_URL}{filename}")

    return Response({
        "status": "success",
        "pptx_file_url": file_url
    })


from django.shortcuts import render
from django.http import HttpResponse

from rest_framework.decorators import api_view
from rest_framework.response import Response
from app1.serializers import SlideHistorySerializer

@api_view(['POST'])
def save_history(request):
    serializer = SlideHistorySerializer(data=request.data)
    if serializer.is_valid():
        serializer.save()
        return Response(serializer.data)
    return Response(serializer.errors, status=400)

from rest_framework.decorators import api_view
from rest_framework.response import Response
from app1.models import SlideHistory
from app1.serializers import SlideHistorySerializer

@api_view(['GET'])
def user_history(request):
    user_id = request.GET.get("user_id")
    data = SlideHistory.objects.filter(user_id=user_id).order_by("-created_at")
    serializer = SlideHistorySerializer(data, many=True)
    return Response(serializer.data)

from django.http import FileResponse
from rest_framework.decorators import api_view
from rest_framework.response import Response
from app1.models import SlideHistory
import os

@api_view(['GET'])
def download_ppt(request, id):
    try:
        record = SlideHistory.objects.get(id=id)
        return FileResponse(record.ppt_file.open(), as_attachment=True, filename=os.path.basename(record.ppt_file.name))
    except SlideHistory.DoesNotExist:
        return Response({"error": "File not found"}, status=404)
